import os
import hashlib
import zipfile
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from ai_engine import AIEngine
from database import DatabaseManager


class FileIndexer:
    SUPPORTED_TEXT = [".pdf", ".docx", ".txt", ".zip", ".ppt", ".pptx"]

    def __init__(self):
        self.ai = AIEngine()
        self.db = DatabaseManager()

    # ======================================================
    # MAIN ENTRY
    # ======================================================
    def process_file(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            self.process_pdf(file_path)

        elif ext == ".docx":
            self.process_docx(file_path)

        elif ext == ".txt":
            self.process_txt(file_path)

        elif ext == ".pptx":
            self.process_pptx_file(file_path)

        elif ext == ".ppt":
            print(f"[Unsupported] Legacy .ppt detected: {os.path.basename(file_path)}")
            print("→ Please convert it to .pptx for indexing.\n")

        elif ext == ".zip":
            self.process_zip(file_path)

    # ======================================================
    # HASH
    # ======================================================
    def compute_hash(self, file_path):
        sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                sha256.update(chunk)
        return sha256.hexdigest()

    # ======================================================
    # COMMON INDEX LOGIC
    # ======================================================
    def index_text_content(self, file_path, full_text):
        if not full_text.strip():
            return

        file_name = os.path.basename(file_path)
        last_modified = os.path.getmtime(file_path)
        file_hash = self.compute_hash(file_path)

        stored_hash = self.db.get_file_hash(file_path)

        if stored_hash and stored_hash == file_hash:
            print(f"Skipped: {file_name}")
            return

        if stored_hash:
            self.db.delete_file_chunks(file_path)

        chunks = self.chunk_text(full_text)

        for chunk in chunks:
            vector = self.ai.text_to_vector(chunk, is_query=False)
            self.db.insert_text_chunk(
                file_path,
                file_name,
                file_hash,
                last_modified,
                chunk,
                vector
            )

        print(f"Indexed: {file_name} ({len(chunks)} chunks)")

    # ======================================================
    # PDF
    # ======================================================
    def process_pdf(self, file_path):
        try:
            reader = PdfReader(file_path)
            full_text = ""

            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    full_text += extracted + "\n"

            self.index_text_content(file_path, full_text)

        except Exception as e:
            print(f"PDF error: {e}")

    # ======================================================
    # DOCX
    # ======================================================
    def process_docx(self, file_path):
        try:
            doc = Document(file_path)
            full_text = "\n".join([para.text for para in doc.paragraphs])
            self.index_text_content(file_path, full_text)
        except Exception as e:
            print(f"DOCX error: {e}")

    # ======================================================
    # TXT
    # ======================================================
    def process_txt(self, file_path):
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                full_text = f.read()
            self.index_text_content(file_path, full_text)
        except Exception as e:
            print(f"TXT error: {e}")

    # ======================================================
    # PPT / PPTX
    # ======================================================
    def process_pptx_file(self, file_path):
        from pptx import Presentation

        try:
            file_name = os.path.basename(file_path)
            last_modified = os.path.getmtime(file_path)
            file_hash = self.compute_file_hash(file_path)

            stored_hash = self.db.get_file_hash(file_path)

            if stored_hash and stored_hash == file_hash:
                print(f"Skipped: {file_name}")
                return

            if stored_hash and stored_hash != file_hash:
                self.db.delete_file_chunks(file_path)

            prs = Presentation(file_path)
            full_text = ""

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            full_text += paragraph.text + "\n"

            if not full_text.strip():
                print(f"No text found in PPTX: {file_name}")
                return

            chunks = self.chunk_text(full_text)

            for chunk in chunks:
                vector = self.ai.text_to_vector(chunk, is_query=False)
                self.db.insert_text_chunk(
                    file_path,
                    file_name,
                    file_hash,
                    last_modified,
                    chunk,
                    vector
                )

            print(f"Indexed: {file_name} ({len(chunks)} chunks)")

        except Exception as e:
            print(f"PPTX error: {file_path} → {e}")
    # ======================================================
    # ZIP (recursive)
    # ======================================================
    def process_zip(self, file_path):
        try:
            with zipfile.ZipFile(file_path, "r") as zip_ref:
                for name in zip_ref.namelist():

                    ext = os.path.splitext(name)[1].lower()

                    if ext == ".txt":

                        with zip_ref.open(name) as f:
                            content = f.read().decode("utf-8", errors="ignore")

                            # Use parent ZIP path for DB tracking
                            virtual_path = f"{file_path}__{name}"

                            self.index_text_content(virtual_path, content)

        except Exception as e:
            print(f"ZIP error: {e}")

    # ======================================================
    # CHUNKING
    # ======================================================
    def chunk_text(self, text, max_length=800):
        paragraphs = text.split("\n")
        chunks = []
        current = ""

        for para in paragraphs:
            if len(current) + len(para) < max_length:
                current += " " + para
            else:
                if current.strip():
                    chunks.append(current.strip())
                current = para

        if current.strip():
            chunks.append(current.strip())

        return chunks